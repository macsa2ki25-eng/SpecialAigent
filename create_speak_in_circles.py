"""
Speak in Circles - Speaking Activity PowerPoint Generator
For high school English classes (CEFR A2-B1, progressive difficulty)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import math

# ── Slide dimensions (16:9) ──
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# ── Colors ──
TEAL = RGBColor(0x00, 0xB4, 0xD8)
PINK = RGBColor(0xFF, 0x6B, 0x8A)
LIGHT_BLUE = RGBColor(0x87, 0xCE, 0xEB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
RED = RGBColor(0xFF, 0x00, 0x00)
BLUE_ACCENT = RGBColor(0x00, 0x7A, 0xCC)
RED_ACCENT = RGBColor(0xFF, 0x44, 0x44)
BG_LIGHT = RGBColor(0xF0, 0xF8, 0xFF)
NAVY = RGBColor(0x1A, 0x23, 0x5B)
GOLD = RGBColor(0xFF, 0xD7, 0x00)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rounded_rect(slide, left, top, width, height, fill_color, border_color=None, border_width=Pt(2)):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 bold=False, color=BLACK, alignment=PP_ALIGN.LEFT,
                 font_name="Calibri", anchor=MSO_ANCHOR.MIDDLE):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    tf.vertical_anchor = anchor
    return txBox


def add_title_bar(slide, set_number):
    """Add the SPEAK in CIRCLES title bar at top."""
    bar = add_rounded_rect(
        slide, Inches(0.3), Inches(0.2), Inches(3.5), Inches(0.7),
        TEAL, border_color=WHITE, border_width=Pt(3)
    )
    tf = bar.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = f"SPEAK in CIRCLES"
    run.font.size = Pt(20)
    run.font.bold = True
    run.font.color.rgb = WHITE
    run.font.name = "Calibri"

    # Set number badge
    badge = add_rounded_rect(
        slide, Inches(4.0), Inches(0.2), Inches(0.9), Inches(0.7),
        NAVY, border_color=WHITE, border_width=Pt(2)
    )
    tf2 = badge.text_frame
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = f"#{set_number}"
    run2.font.size = Pt(22)
    run2.font.bold = True
    run2.font.color.rgb = WHITE
    run2.font.name = "Calibri"


def create_title_slide(prs):
    """Create the title/instructions slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    set_slide_bg(slide, BG_LIGHT)

    # Main title
    add_text_box(slide, Inches(1.5), Inches(0.5), Inches(10), Inches(1.5),
                 "SPEAK in CIRCLES", font_size=54, bold=True, color=TEAL,
                 alignment=PP_ALIGN.CENTER, font_name="Calibri")

    # Subtitle
    add_text_box(slide, Inches(2), Inches(1.8), Inches(9), Inches(0.6),
                 "Say it without saying it!", font_size=28, bold=False,
                 color=PINK, alignment=PP_ALIGN.CENTER, font_name="Calibri")

    # Decorative line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(3), Inches(2.5), Inches(7), Pt(3)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = TEAL
    line.line.fill.background()

    # Instructions box
    instr_box = add_rounded_rect(
        slide, Inches(1.5), Inches(2.8), Inches(10), Inches(4.2),
        WHITE, border_color=TEAL, border_width=Pt(3)
    )

    # Step 1
    step1_badge = add_rounded_rect(
        slide, Inches(1.8), Inches(3.1), Inches(0.7), Inches(0.7),
        TEAL
    )
    tf = step1_badge.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "1"
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = WHITE

    add_text_box(slide, Inches(2.8), Inches(3.0), Inches(8), Inches(0.9),
                 "Student A sees the picture and memorizes the object.\nStudent B closes their eyes.",
                 font_size=18, color=DARK_GRAY, alignment=PP_ALIGN.LEFT)

    # Step 2
    step2_badge = add_rounded_rect(
        slide, Inches(1.8), Inches(4.1), Inches(0.7), Inches(0.7),
        PINK
    )
    tf = step2_badge.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "2"
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = WHITE

    add_text_box(slide, Inches(2.8), Inches(4.0), Inches(8), Inches(0.9),
                 "Student B opens their eyes.\nStudent A gives hints without saying the name.",
                 font_size=18, color=DARK_GRAY, alignment=PP_ALIGN.LEFT)

    # Step 3
    step3_badge = add_rounded_rect(
        slide, Inches(1.8), Inches(5.2), Inches(0.7), Inches(0.7),
        LIGHT_BLUE
    )
    tf = step3_badge.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "3"
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = WHITE

    add_text_box(slide, Inches(2.8), Inches(5.1), Inches(8), Inches(0.9),
                 "Choices appear.\nIf Student B picks the right one — Mission Complete!",
                 font_size=18, color=DARK_GRAY, alignment=PP_ALIGN.LEFT)

    # Footer
    add_text_box(slide, Inches(2), Inches(6.5), Inches(9), Inches(0.5),
                 "Speaking Activity  |  CEFR A2–B1 (Progressive)  |  10 Sets",
                 font_size=14, color=TEAL, alignment=PP_ALIGN.CENTER)


def create_photo_slide(prs, set_number, item_name, photo_description, level="A2"):
    """Create the photo slide with description text (no actual photo)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_LIGHT)
    add_title_bar(slide, set_number)

    # Level badge
    level_color = TEAL if level == "A2" else PINK if level == "A2-B1" else RGBColor(0xFF, 0x45, 0x00)
    level_box = add_rounded_rect(
        slide, Inches(5.2), Inches(0.2), Inches(1.5), Inches(0.7),
        level_color
    )
    tf = level_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = level
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.color.rgb = WHITE

    # Phase indicator
    phase_box = add_rounded_rect(
        slide, Inches(9.5), Inches(0.2), Inches(3.5), Inches(0.7),
        PINK
    )
    tf = phase_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "👀 Student A: Look & Memorize!"
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = WHITE

    # Large photo placeholder area
    photo_area = add_rounded_rect(
        slide, Inches(1.5), Inches(1.3), Inches(10.3), Inches(5.2),
        WHITE, border_color=TEAL, border_width=Pt(4)
    )

    # Photo description (since we can't embed actual photos)
    add_text_box(slide, Inches(2), Inches(1.8), Inches(9.3), Inches(1.0),
                 f"📷 Photo: {photo_description}",
                 font_size=24, bold=True, color=DARK_GRAY,
                 alignment=PP_ALIGN.CENTER)

    # Large icon/emoji representation
    add_text_box(slide, Inches(3), Inches(2.8), Inches(7), Inches(2.5),
                 get_emoji_for_item(item_name),
                 font_size=120, bold=False, color=BLACK,
                 alignment=PP_ALIGN.CENTER)

    # Instruction text
    add_text_box(slide, Inches(2), Inches(5.5), Inches(9), Inches(0.8),
                 "Student B: Please close your eyes!",
                 font_size=20, bold=True, color=PINK,
                 alignment=PP_ALIGN.CENTER)

    # Timer hint
    timer_badge = add_rounded_rect(
        slide, Inches(10.5), Inches(5.8), Inches(1.5), Inches(0.8),
        NAVY
    )
    tf = timer_badge.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "7 sec"
    run.font.size = Pt(20)
    run.font.bold = True
    run.font.color.rgb = WHITE


def create_choices_slide(prs, set_number, correct_answer, choice_a, choice_b, choice_c):
    """Create the 3-choice slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_LIGHT)
    add_title_bar(slide, set_number)

    # Phase indicator
    phase_box = add_rounded_rect(
        slide, Inches(9.5), Inches(0.2), Inches(3.5), Inches(0.7),
        TEAL
    )
    tf = phase_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "🤔 Student B: Choose the answer!"
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = WHITE

    choices = [
        ("A", choice_a, TEAL),
        ("B", choice_b, PINK),
        ("C", choice_c, LIGHT_BLUE),
    ]

    y_start = Inches(1.5)
    for i, (letter, text, color) in enumerate(choices):
        y = y_start + Inches(i * 1.8)

        # Choice button
        btn = add_rounded_rect(
            slide, Inches(2), y, Inches(9), Inches(1.4),
            color, border_color=WHITE, border_width=Pt(4)
        )

        # Letter badge
        letter_box = add_rounded_rect(
            slide, Inches(2.3), y + Inches(0.2), Inches(1.0), Inches(1.0),
            WHITE, border_color=WHITE, border_width=Pt(2)
        )
        tf = letter_box.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = letter
        run.font.size = Pt(36)
        run.font.bold = True
        run.font.color.rgb = color

        # Choice text
        add_text_box(slide, Inches(3.8), y, Inches(6.5), Inches(1.4),
                     text, font_size=40, bold=True, color=WHITE,
                     alignment=PP_ALIGN.CENTER)

    # Note at bottom
    add_text_box(slide, Inches(2), Inches(6.8), Inches(9), Inches(0.5),
                 f"✅ Correct Answer: {correct_answer}",
                 font_size=14, bold=False, color=RGBColor(0x99, 0x99, 0x99),
                 alignment=PP_ALIGN.RIGHT)


def create_model_slide(prs, set_number, item_name, vocab_list, model_sentence, grammar_point=""):
    """Create the model description slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_LIGHT)
    add_title_bar(slide, set_number)

    # Phase indicator
    phase_box = add_rounded_rect(
        slide, Inches(7), Inches(0.2), Inches(6), Inches(0.7),
        GOLD
    )
    tf = phase_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "⭐ Model Description"
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.color.rgb = DARK_GRAY

    # Grammar / Vocabulary box
    vocab_box = add_rounded_rect(
        slide, Inches(0.5), Inches(1.2), Inches(12.3), Inches(2.8),
        WHITE, border_color=DARK_GRAY, border_width=Pt(2)
    )

    # Grammar/Vocabulary header
    header_box = add_rounded_rect(
        slide, Inches(4.5), Inches(1.0), Inches(4.3), Inches(0.5),
        DARK_GRAY
    )
    tf = header_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Grammar / Vocabulary"
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = WHITE

    # Answer label
    add_text_box(slide, Inches(1.0), Inches(1.6), Inches(4), Inches(0.6),
                 f"Answer: {item_name}",
                 font_size=22, bold=True, color=TEAL,
                 alignment=PP_ALIGN.LEFT)

    # Vocabulary items
    y_offset = Inches(2.2)
    colors = [BLUE_ACCENT, RED_ACCENT, TEAL, PINK]
    for i, vocab in enumerate(vocab_list):
        color = colors[i % len(colors)]
        # Bullet
        bullet = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(1.2), y_offset + Inches(i * 0.5) + Inches(0.1),
            Inches(0.2), Inches(0.2)
        )
        bullet.fill.solid()
        bullet.fill.fore_color.rgb = color
        bullet.line.fill.background()

        add_text_box(slide, Inches(1.6), y_offset + Inches(i * 0.5), Inches(10), Inches(0.5),
                     vocab, font_size=18, bold=True, color=DARK_GRAY,
                     alignment=PP_ALIGN.LEFT)

    # Grammar point (if any)
    if grammar_point:
        gp_box = add_rounded_rect(
            slide, Inches(7), Inches(2.2), Inches(5.5), Inches(0.5),
            RGBColor(0xE8, 0xE8, 0xE8)
        )
        tf = gp_box.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = f"📝 {grammar_point}"
        run.font.size = Pt(14)
        run.font.color.rgb = DARK_GRAY

    # Divider line
    divider = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(4.3), Inches(11.7), Pt(2)
    )
    divider.fill.solid()
    divider.fill.fore_color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    divider.line.fill.background()

    # Model sentence box
    sentence_box = add_rounded_rect(
        slide, Inches(0.5), Inches(4.6), Inches(12.3), Inches(2.5),
        WHITE, border_color=TEAL, border_width=Pt(3)
    )

    add_text_box(slide, Inches(1.0), Inches(4.8), Inches(11.3), Inches(2.0),
                 model_sentence,
                 font_size=28, bold=True, color=DARK_GRAY,
                 alignment=PP_ALIGN.CENTER, font_name="Calibri")


def get_emoji_for_item(item_name):
    """Return a relevant emoji for the item."""
    emoji_map = {
        "Bicycle": "🚲",
        "Umbrella": "☂️",
        "Violin": "🎻",
        "Sunglasses": "🕶️",
        "Pancake": "🥞",
        "Microscope": "🔬",
        "Thermometer": "🌡️",
        "Sleeping Bag": "🛏️",
        "Yogurt": "🥛",
        "Submarine": "🚢",
    }
    return emoji_map.get(item_name, "❓")


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # ── Activity Data: 10 Sets (Progressive: A2 → A2-B1 → B1) ──
    sets = [
        # ===== Sets 1-4: CEFR A2 (Basic, everyday objects, simple descriptions) =====
        {
            "number": 1,
            "level": "A2",
            "item": "Bicycle",
            "photo_desc": "A woman riding a bicycle on a city street near a brick bridge",
            "choices": ("Bicycle", "Motorcycle", "Scooter"),
            "correct": "A",
            "vocab": ["pedal — ペダル（足でこぐ部分）", "two wheels — 2つの車輪"],
            "model": "It has two wheels. You push the pedals with your feet to move. It does not need gasoline.",
            "grammar": "Comparative: faster than ~",
        },
        {
            "number": 2,
            "level": "A2",
            "item": "Umbrella",
            "photo_desc": "A person holding a colorful umbrella on a rainy street",
            "choices": ("Parasol", "Umbrella", "Raincoat"),
            "correct": "B",
            "vocab": ["fold — 折りたたむ", "above your head — 頭の上に"],
            "model": "You hold it above your head when it rains. You can fold it and put it in your bag.",
            "grammar": "When + present tense",
        },
        {
            "number": 3,
            "level": "A2",
            "item": "Violin",
            "photo_desc": "A musician playing the violin on stage",
            "choices": ("Cello", "Violin", "Guitar"),
            "correct": "B",
            "vocab": ["strings — 弦", "bow — 弓（弦楽器を弾く道具）"],
            "model": "It has strings and you use a bow to play it. You hold it under your chin. It is smaller than a cello.",
            "grammar": "Subject + have/has + noun",
        },
        {
            "number": 4,
            "level": "A2",
            "item": "Sunglasses",
            "photo_desc": "A person wearing sunglasses at the beach on a bright sunny day",
            "choices": ("Goggles", "Sunglasses", "Reading Glasses"),
            "correct": "B",
            "vocab": ["lenses — レンズ", "shade — 日よけ、色の濃い部分"],
            "model": "You wear them on your face outdoors. They have dark lenses that protect your eyes from the sun.",
            "grammar": "Relative pronoun: that",
        },
        # ===== Sets 5-7: CEFR A2-B1 (Transitional, more abstract, longer descriptions) =====
        {
            "number": 5,
            "level": "A2-B1",
            "item": "Pancake",
            "photo_desc": "A stack of fluffy pancakes with syrup and butter on a plate",
            "choices": ("Crepe", "Pancake", "Waffle"),
            "correct": "B",
            "vocab": ["batter — 生地（液状の）", "fluffy — ふわふわの", "pour — 注ぐ"],
            "model": "You make it by pouring batter on a hot pan. It is round and flat, but thicker and fluffier than a crepe. People often eat it for breakfast with syrup.",
            "grammar": "Comparative: thicker / fluffier than ~",
        },
        {
            "number": 6,
            "level": "A2-B1",
            "item": "Microscope",
            "photo_desc": "A scientist looking through a microscope in a laboratory",
            "choices": ("Telescope", "Microscope", "Binoculars"),
            "correct": "B",
            "vocab": ["magnify — 拡大する", "lens — レンズ", "tiny — とても小さい"],
            "model": "It is used in science class to look at very tiny things. It magnifies objects that are too small for your eyes. Unlike a telescope, you look down into it.",
            "grammar": "Passive: It is used to ~",
        },
        {
            "number": 7,
            "level": "A2-B1",
            "item": "Thermometer",
            "photo_desc": "A digital thermometer showing a person's body temperature",
            "choices": ("Thermometer", "Barometer", "Stopwatch"),
            "correct": "A",
            "vocab": ["temperature — 温度", "measure — 測る", "fever — 熱"],
            "model": "You use it to measure how hot or cold something is. When you are sick, your mother may put it in your mouth to check if you have a fever.",
            "grammar": "Interrogative: how + adjective",
        },
        # ===== Sets 8-10: CEFR B1 (More complex, abstract concepts, nuanced descriptions) =====
        {
            "number": 8,
            "level": "B1",
            "item": "Sleeping Bag",
            "photo_desc": "A person inside a sleeping bag at a campsite in the mountains",
            "choices": ("Blanket", "Sleeping Bag", "Tent"),
            "correct": "B",
            "vocab": ["zipper — ジッパー", "insulated — 断熱された", "portable — 持ち運びできる"],
            "model": "It is a portable item you use when you camp outdoors. You get inside it to stay warm at night. It has a zipper and you can roll it up, unlike a blanket which has no zipper.",
            "grammar": "Relative clause: which + verb",
        },
        {
            "number": 9,
            "level": "B1",
            "item": "Yogurt",
            "photo_desc": "A cup of yogurt with fruit toppings on a table",
            "choices": ("Pudding", "Yogurt", "Ice Cream"),
            "correct": "B",
            "vocab": ["fermented — 発酵した", "bacteria — 細菌、バクテリア", "sour — 酸っぱい"],
            "model": "It is a dairy product made from fermented milk. It has a slightly sour taste and a creamy texture. Unlike ice cream, it is not frozen, and unlike pudding, it is made with bacteria that are good for your health.",
            "grammar": "Contrast: unlike ~ / whereas ~",
        },
        {
            "number": 10,
            "level": "B1",
            "item": "Submarine",
            "photo_desc": "A submarine diving deep under the ocean surface",
            "choices": ("Ferry", "Submarine", "Cruise Ship"),
            "correct": "B",
            "vocab": ["dive — 潜る", "surface — 水面", "crew — 乗組員", "underwater — 水中の"],
            "model": "It is a type of vessel that is designed to travel underwater. It can dive below the surface and stay there for a long time. The crew cannot open the windows because the water pressure outside is extremely high.",
            "grammar": "Reason clause: because ~",
        },
    ]

    # ── Create Title Slide ──
    create_title_slide(prs)

    # ── Create 10 Sets (3 slides each) ──
    for s in sets:
        # Slide 1: Photo
        create_photo_slide(prs, s["number"], s["item"], s["photo_desc"], s.get("level", "A2"))

        # Slide 2: Choices
        create_choices_slide(
            prs, s["number"], s["correct"],
            s["choices"][0], s["choices"][1], s["choices"][2]
        )

        # Slide 3: Model Description
        create_model_slide(
            prs, s["number"], s["item"],
            s["vocab"], s["model"], s["grammar"]
        )

    # ── Save ──
    output_path = "/home/user/SpecialAigent/Speak_in_Circles_Activity.pptx"
    prs.save(output_path)
    print(f"✅ PowerPoint saved to: {output_path}")
    print(f"📊 Total slides: {len(prs.slides)} (1 title + 10 sets × 3 slides)")


if __name__ == "__main__":
    main()
